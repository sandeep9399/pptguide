import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import io
import requests

st.set_page_config(page_title="Apollo PPT Converter", layout="wide")
st.title("🧩 Convert Old PPT to Apollo Standard")

st.markdown("""
Upload your **old-style PowerPoint presentation (.pptx)** below. This tool will:
- Extract your content slide-by-slide
- Reformat using Apollo University styles (title, footer, clean layout)
- Add branding: Apollo University logo and Powered by Apollo Knowledge
- Apply light blue theme to all slides
- Output a new `.pptx` file ready to download
""")

uploaded_file = st.file_uploader("Upload your old-format PPT", type=["pptx"])
apollo_logo_url = "https://upload.wikimedia.org/wikipedia/en/1/1e/Apollo_Hospitals_Logo.png"

def suggest_design_elements(full_text):
    full_text = full_text.lower()
    if "who" in full_text:
        return {"Layout": "Two-column: WHO quote on left, image on right", "Visual": "Quote bubble + doctor team image"}
    elif "components" in full_text:
        return {"Layout": "4-quadrant grid", "Visual": "Infographic with icons"}
    elif "india" in full_text:
        return {"Layout": "Map overlay", "Visual": "India map infographic"}
    else:
        return {"Layout": "Standard title-content", "Visual": "Photo + icon"}

if uploaded_file:
    old_ppt = Presentation(uploaded_file)
    new_ppt = Presentation()
    layout = new_ppt.slide_layouts[1]

    for i, old_slide in enumerate(old_ppt.slides, start=1):
        new_slide = new_ppt.slides.add_slide(layout)
        try:
            old_title = old_slide.shapes.title.text.strip()
            new_slide.shapes.title.text = old_title
            title_para = new_slide.shapes.title.text_frame.paragraphs[0]
            title_para.font.name = "Poppins"
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(0, 51, 102)
        except:
            new_slide.shapes.title.text = f"Slide {i}"

        text_content = []
        for shape in old_slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if text: text_content.append(text)

        content_box = new_slide.placeholders[1].text_frame
        content_box.clear()
        for line in text_content:
            para = content_box.add_paragraph()
            para.text = line
            para.font.size = Pt(18)
            para.font.name = "Segoe UI"
            para.font.color.rgb = RGBColor(0, 0, 0)

        combined_text = " ".join(text_content)
        suggestion = suggest_design_elements(combined_text)
        layout_box = new_slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(5.5), Inches(0.6))
        layout_tf = layout_box.text_frame
        layout_tf.text = f"AI Layout: {suggestion['Layout']} | Visual: {suggestion['Visual']}"

        footer_box = new_slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.3))
        footer_tf = footer_box.text_frame
        footer_tf.text = "Powered by Apollo Knowledge"
        footer_para = footer_tf.paragraphs[0]
        footer_para.font.name = "Segoe UI"
        footer_para.font.size = Pt(10)
        footer_para.font.italic = True
        footer_para.font.color.rgb = RGBColor(100, 100, 100)

        try:
            response = requests.get(apollo_logo_url)
            if response.status_code == 200:
                img_data = io.BytesIO(response.content)
                new_slide.shapes.add_picture(img_data, Inches(8.2), Inches(6.7), width=Inches(1))
        except: pass

        fill = new_slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(210, 230, 255)

    ppt_io = io.BytesIO()
    new_ppt.save(ppt_io)
    ppt_io.seek(0)
    st.download_button("🎯 Download AI-Enhanced Apollo Slides", data=ppt_io, file_name="Apollo_Enhanced_Presentation.pptx")