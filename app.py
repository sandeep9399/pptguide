# Apollo Slide Visual Guide App (v2.2)
# Upload a PPT → Get design suggestions + inline image previews for visual prompts

import streamlit as st
import pandas as pd
from pptx import Presentation
import io
import re
import urllib.parse

st.set_page_config(page_title="PPT Visual Guide Generator", layout="wide")
st.title("🎯 Apollo PPT Visual Design Guide")
st.markdown("Upload your PowerPoint presentation below. We'll analyze and generate layout, visual suggestions, image ideas, and export-ready design guides.")

uploaded_file = st.file_uploader("Upload your .pptx file", type=["pptx"])

def clean_illegal_chars(text):
    ILLEGAL_CHARACTERS_RE = re.compile(r"[\000-\010\013\014\016-\037]")
    return ILLEGAL_CHARACTERS_RE.sub("", text)

def suggest_design_elements(full_text):
    if "who" in full_text:
        return {
            "Suggested Layout": "Two-column: WHO quote on left, image on right",
            "Typography": "Poppins Bold 36pt title, Segoe UI 22pt body",
            "Color Theme": "Blue-Grey healthcare theme",
            "Icon Style": "Line icons (globe, stethoscope)",
            "Animation": "Fade-in for text, fly-in for image",
            "Visual Type": "Image of WHO HQ or healthcare team + quote callout",
            "Visual Prompt": "High-resolution image of WHO healthcare theme, quote bubble left, global doctor team on right"
        }
    elif "components" in full_text:
        return {
            "Suggested Layout": "4-quadrant grid",
            "Typography": "Arial Rounded 28pt bold titles, 20pt text",
            "Color Theme": "Color blocks: blue, green, orange, violet",
            "Icon Style": "Health category icons (🧠 ❤️ 🧘‍♀️ 👥)",
            "Animation": "Sequential fade-in",
            "Visual Type": "Infographic grid with physical, mental, social, spiritual labels",
            "Visual Prompt": "Flat-style infographic with 4 quadrants labeled Physical, Mental, Social, Spiritual using healthcare icons"
        }
    elif "india" in full_text:
        return {
            "Suggested Layout": "Data + map overlay",
            "Typography": "Segoe UI Bold 30pt, Regular 20pt",
            "Color Theme": "Warm tones + India map overlay",
            "Icon Style": "Flat infographic symbols",
            "Animation": "Bar chart build-up",
            "Visual Type": "Map of India with NCD statistics, hotspot callouts",
            "Visual Prompt": "Infographic map of India showing non-communicable disease hotspots and stats with health icons"
        }
    else:
        return {
            "Suggested Layout": "Title + Image Right",
            "Typography": "Calibri Light 32pt title, 20pt body",
            "Color Theme": "Light pastel healthcare theme",
            "Icon Style": "Simple flat icons",
            "Animation": "Appear on click",
            "Visual Type": "Healthcare teamwork photo + soft icons",
            "Visual Prompt": "Soft pastel-themed slide with doctor team photo on right and clean healthcare icons"
        }

def analyze_ppt(ppt_file):
    prs = Presentation(ppt_file)
    data = []

    for i, slide in enumerate(prs.slides, start=1):
        block_texts = [
            clean_illegal_chars(shape.text.strip()) for shape in slide.shapes
            if shape.has_text_frame and shape.text.strip()
        ]
        full_text = " ".join(block_texts).lower()
        split_required = len(block_texts) > 3 or any(k in full_text for k in ["components", "diseases", "determinants", "definition"])

        part_count = 2 if split_required else 1
        for idx in range(part_count):
            part = f"Part {idx+1}" if part_count == 2 else "Full Slide"
            design = suggest_design_elements(full_text)

            search_term = design['Visual Prompt'].split(',')[0].split(' ')[:5]
            preview_prompt_url = f"https://source.unsplash.com/800x400/?{urllib.parse.quote(' '.join(search_term))}"

            data.append({
                "Slide Number": i,
                "Slide Part": part,
                "Block Title": block_texts[0][:60] + "..." if block_texts else "Untitled",
                "Content Alignment": design["Suggested Layout"],
                "Suggested Visual Style": design["Icon Style"],
                "Designer Note": "Split this into multiple sections." if part != "Full Slide" else "All content fits in one slide.",
                "Typography": design["Typography"],
                "Color Theme": design["Color Theme"],
                "Animation": design["Animation"],
                "Visual Type": design["Visual Type"],
                "Visual Prompt": design["Visual Prompt"],
                "Preview Image URL": preview_prompt_url
            })

    return pd.DataFrame(data)

if uploaded_file:
    df = analyze_ppt(uploaded_file)
    st.success("✅ Analysis complete!")
    st.dataframe(df.drop(columns=["Preview Image URL"]), use_container_width=True)

    # Show visual preview thumbnails
    st.markdown("### 🖼️ Visual Previews from Prompts")
    for _, row in df.iterrows():
        st.markdown(f"**Slide {row['Slide Number']} - {row['Slide Part']}**")
        st.image(row["Preview Image URL"], caption=row["Visual Prompt"], use_container_width=True)

    # Excel export
    towrite = io.BytesIO()
    df.drop(columns=["Preview Image URL"]).to_excel(towrite, index=False, engine='openpyxl')
    towrite.seek(0)

    st.download_button(
        label="📥 Download Excel",
        data=towrite,
        file_name="Slide_Visual_Design_Guide.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )
